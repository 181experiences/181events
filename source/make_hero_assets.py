#!/usr/bin/env python3
"""Renders each event's web hero as real files, ready for the Assets screen.

For every distinct event design in the calendar, two files land in ../print/heroes/:
  {slug}_web-hero.jpg           1600 x 900, exactly what the site shows: the
                                gradient, or the typographic card with its title
  {slug}_web-hero-editable.pptx the same artwork as a PowerPoint: background
                                image with the title as a live text box. Canva
                                imports .pptx with text still editable, so this
                                is the working file for restyling.

Fonts follow make_qr_signs.py: Hanken Grotesk from the brand kit, Georgia as
the Marcellus stand-in. Run: python source/make_hero_assets.py"""

import os, re, math
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from events_data import EVENTS

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.normpath(os.path.join(HERE, "..", "print", "heroes"))
os.makedirs(OUT, exist_ok=True)

W, H = 1600, 900
INK = (22, 22, 26); PAPER = (247, 244, 239); RED = (196, 31, 38); STONE = (168, 160, 148)

def _find_hanken():
    d = HERE
    for _ in range(10):
        p = os.path.join(d, "tidemere", "2_Marketing", "brand", "fonts", "HankenGrotesk-Variable.ttf")
        if os.path.exists(p):
            return p
        nd = os.path.dirname(d)
        if nd == d:
            break
        d = nd
    return "C:/Windows/Fonts/segoeui.ttf"

HANKEN = _find_hanken()
MARCELLUS = os.path.join(HERE, "marcellus.ttf")
GEORGIA = "C:/Windows/Fonts/georgia.ttf"

def display_font(size):
    return ImageFont.truetype(MARCELLUS if os.path.exists(MARCELLUS) else GEORGIA, size)

def body_font(size, weight=500):
    f = ImageFont.truetype(HANKEN, size)
    try: f.set_variation_by_axes([weight])
    except Exception: pass
    return f

def plain(s):
    for a, b in [("&rsquo;", "\u2019"), ("&amp;", "&"), ("&eacute;", "\u00e9"),
                 ("<em>", ""), ("</em>", ""), ("&mdash;", "-")]:
        s = s.replace(a, b)
    return s

def parse_gradient(css):
    m = re.match(r"linear-gradient\((\d+)deg,(.+)\)", css)
    angle = int(m.group(1))
    stops = []
    for part in m.group(2).split(","):
        cm = re.search(r"#([0-9a-fA-F]{6})\s+(\d+)%", part)
        stops.append((int(cm.group(2)) / 100.0,
                      tuple(int(cm.group(1)[i:i+2], 16) for i in (0, 2, 4))))
    return angle, stops

def gradient_image(css):
    """CSS linear-gradient -> 1600x900. A rotated luminance ramp indexes a lookup
    table built from the stops; close enough to the browser to pass for it."""
    angle, stops = parse_gradient(css)
    big = 2400
    ramp = Image.linear_gradient("L").resize((big, big))          # 0 top -> 255 bottom
    ramp = ramp.rotate(180 - angle, resample=Image.BILINEAR)      # CSS angles run clockwise from up
    ramp = ramp.crop((big//2 - W//2, big//2 - H//2, big//2 + W//2, big//2 + H//2))
    luts = [[0]*256 for _ in range(3)]
    for v in range(256):
        t = v / 255.0
        lo = stops[0]
        hi = stops[-1]
        for i in range(len(stops) - 1):
            if stops[i][0] <= t <= stops[i+1][0]:
                lo, hi = stops[i], stops[i+1]
                break
        span = max(hi[0] - lo[0], 1e-6)
        f = min(1.0, max(0.0, (t - lo[0]) / span))
        for c in range(3):
            luts[c][v] = round(lo[1][c] + (hi[1][c] - lo[1][c]) * f)
    bands = [ramp.point(luts[c]) for c in range(3)]
    return Image.merge("RGB", bands)

def tracked(draw, cx, y, text, font, tracking, fill):
    widths = [draw.textlength(c, font=font) for c in text]
    total = sum(widths) + tracking * (len(text) - 1)
    x = cx - total / 2
    for c, w in zip(text, widths):
        draw.text((x, y), c, font=font, fill=fill)
        x += w + tracking

def card_base():
    """The typographic card's stage: ink field and the red rule, no text."""
    img = Image.new("RGB", (W, H), INK)
    d = ImageDraw.Draw(img)
    d.rectangle([W/2 - 44, 300, W/2 + 44, 304], fill=RED)
    return img

def card_with_text(title):
    img = card_base()
    d = ImageDraw.Draw(img)
    size = 88 if len(title) <= 26 else 72 if len(title) <= 36 else 58
    f = display_font(size)
    w = d.textlength(title, font=f)
    d.text((W/2 - w/2, 400), title, font=f, fill=PAPER)
    tracked(d, W/2, 400 + size + 62, "181 FREMONT", body_font(26, 500), 9, STONE)
    return img

def editable_pptx(bg_path, title, out_path, dark_text=False):
    prs = Presentation()
    prs.slide_width = Emu(int(Inches(13.333)))
    prs.slide_height = Emu(int(Inches(7.5)))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    run = para.add_run(); run.text = title
    run.font.name = "Georgia"; run.font.size = Pt(48)
    run.font.color.rgb = RGBColor(0x16, 0x16, 0x1A) if dark_text else RGBColor(0xF7, 0xF4, 0xEF)
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.6))
    sp = sub.text_frame.paragraphs[0]; sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run(); sr.text = "1 8 1   F R E M O N T"
    sr.font.name = "Hanken Grotesk"; sr.font.size = Pt(14)
    sr.font.color.rgb = RGBColor(0xA8, 0xA0, 0x94)
    prs.save(out_path)

seen = set()
made = 0
for e in EVENTS:
    slug = e["slug"]
    if slug in seen:
        continue
    seen.add(slug)
    title = plain(e["title"])
    jpg = os.path.join(OUT, f"{slug}_web-hero.jpg")
    if e["img"]:
        img = gradient_image(e["img"])
        img.save(jpg, quality=88)
        editable_pptx(jpg, title, os.path.join(OUT, f"{slug}_web-hero-editable.pptx"))
    else:
        card_with_text(title).save(jpg, quality=90)
        base = os.path.join(OUT, f"_{slug}_stage.jpg")
        card_base().save(base, quality=90)
        editable_pptx(base, title, os.path.join(OUT, f"{slug}_web-hero-editable.pptx"))
        os.remove(base)
    made += 1
print(f"{made} designs -> {OUT}")
